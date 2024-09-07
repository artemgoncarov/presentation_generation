from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches
from random import choice
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import nltk
from nltk.tokenize import sent_tokenize
nltk.download('punkt')

# with open('static/.json', 'r') as file:
#     data = json.load(file)

# name = data['name']
# design = data['design']
# plan = data['plan']
# slides = data['slides']

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
   
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))

def remove_slides(presentation, slide_indices):
    """
    Удаляет слайды по указанным индексам из презентации.

    :param presentation: Объект Presentation
    :param slide_indices: Список индексов слайдов для удаления
    """
    # Получаем XML-узел, содержащий идентификаторы слайдов
    slide_id_list = presentation.slides._sldIdLst
    slide_ids = list(slide_id_list)

    # Удаляем слайды в обратном порядке
    for index in sorted(slide_indices, reverse=True):
        slide_id = slide_ids[index]
        slide_id_list.remove(slide_id)


def generate_presentation_maket(slides, plan, design, name, output_file, font_name='Corbel Light', font_size=26, font_color = (0,0,0), maket="1_шаблон.pptx"):

    presentation = Presentation(maket)

    slides_with_picture = []
    slides_for_text = []
    slides_for_text_small = []
    slides_for_text_huge = []

    count_slides_maket = 0

    for slide_number, slide in enumerate(presentation.slides, ):
        count_slides_maket+=1
        flag=True
        count_for_slides_text = 0
        for placeholder in slide.placeholders:
            if 'Рисунок' in placeholder.name or 'Picture' in placeholder.name:
                  slides_with_picture.append(slide_number)
                  flag = False
                  break
            if 'Текст' in placeholder.name or 'Text' in placeholder.name:
                  count_for_slides_text+=1

        if flag:
            slides_for_text.append(slide_number)
            if count_for_slides_text == 1:
                slides_for_text_small.append(slide_number)
            else:
                slides_for_text_huge.append(slide_number)
    print(slides_for_text_huge)
    if design=='1':
        if 5 in slides_with_picture:
            slides_with_picture.remove(5)
        slides_for_text.remove(4)

    slides_for_text.remove(0)
    if 0 in slides_for_text_huge:
      slides_for_text_huge.remove(0)
    else:
      slides_for_text_small.remove(0)

    slide_layout = presentation.slide_layouts[0]  # Макет для заголовка и содержимого
    slide = presentation.slides.add_slide(slide_layout)

    # Добавляем заголовок
    title = slide.shapes.title
    title.text = name

    # Настраиваем шрифт заголовка
    title.text_frame.paragraphs[0].font.name = font_name
    title.text_frame.paragraphs[0].font.size = Pt(font_size + 20)
    title.text_frame.paragraphs[0].font.bold = True

    count = 0
    # Проходим по каждому элементу содержания (тематика и текст для слайда)
    for slide_content, title_content in zip(slides, plan):

         #считает порядок слайдов
        if count ==len(slides_for_text)-1:
            count = 0
        else:
            count+=1

        if slide_content['image_path'] != '':
            slide_layout = presentation.slide_layouts[choice(slides_with_picture)]  # Макет для заголовка и содержимого
            slide = presentation.slides.add_slide(slide_layout)

            # Добавляем заголовок
            title = slide.shapes.title
            title.text = title_content

            # Настраиваем шрифт заголовка
            title.text_frame.paragraphs[0].font.name = font_name
            title.text_frame.paragraphs[0].font.size = Pt(font_size + 10)
            title.text_frame.paragraphs[0].font.bold = True

            text_forms = []
            img_forms = []
            for placeholder in slide.placeholders:
                if 'Text' in placeholder.name:
                    text_forms.append(placeholder.placeholder_format.idx)
                if 'Picture' in placeholder.name:
                    img_forms.append(placeholder.placeholder_format.idx)

                print(f"Индекс: {placeholder.placeholder_format.idx}, Тип: {placeholder.name}")
            print()

            content_shape = slide.shapes.placeholders[text_forms[0]]
            content_shape.text = slide_content['text']

            content_shape2 = slide.shapes.placeholders[img_forms[-1]]
            content_shape2.insert_picture(slide_content['image_path'])

            # Настраиваем шрифт основного текста
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = RGBColor(*font_color)
                # paragraph.alignment = PP_ALIGN.LEFT

        else:
            if '[r]' not in slide_content['text']:
                try:
                    text_predlozenia = sent_tokenize(slide_content['text'],)
                except:
                    text_predlozenia = list(map(str, slide_content['text'].split('.')))

            else:
                try:
                    text_predlozenia = list(map(str, slide_content['text'].split('[r]')))

                except:
                    try:
                        text_predlozenia = sent_tokenize(slide_content['text'],)
                    except:
                        text_predlozenia = list(map(str, slide_content['text'].split('.')))

            if len(text_predlozenia)<=3:
                try:
                    slide_layout = presentation.slide_layouts[choice(slides_for_text_small)]  # Макет для заголовка и содержимого
                    slide = presentation.slides.add_slide(slide_layout)
                except:
                    try:
                        slide_layout = presentation.slide_layouts[slides_for_text_small[0]]  # Макет для заголовка и содержимого
                        slide = presentation.slides.add_slide(slide_layout)
                    except:
                        slide_layout = presentation.slide_layouts[slides_for_text[0]]  # Макет для заголовка и содержимого
                        slide = presentation.slides.add_slide(slide_layout)
            else:
                print(slides_for_text_huge)
                try:
                    slide_layout = presentation.slide_layouts[choice(slides_for_text_huge)]  # Макет для заголовка и содержимого
                    slide = presentation.slides.add_slide(slide_layout)
                except:
                    try:
                        slide_layout = presentation.slide_layouts[slides_for_text_huge[0]]  # Макет для заголовка и содержимого
                        slide = presentation.slides.add_slide(slide_layout)
                    except:
                        slide_layout = presentation.slide_layouts[slides_for_text[0]]  # Макет для заголовка и содержимого
                        slide = presentation.slides.add_slide(slide_layout)

            # Добавляем заголовок
            title = slide.shapes.title
            title.text = title_content

            # Настраиваем шрифт заголовка
            title.text_frame.paragraphs[0].font.name = font_name
            title.text_frame.paragraphs[0].font.size = Pt(font_size + 10)
            title.text_frame.paragraphs[0].font.bold = True

            text_forms = []
            img_forms = []

            for placeholder in slide.placeholders:
                if 'Text' in placeholder.name:
                    text_forms.append(placeholder.placeholder_format.idx)
                if 'Picture' in placeholder.name:
                    img_forms.append(placeholder.placeholder_format.idx)

            #     print(f"Индекс: {placeholder.placeholder_format.idx}, Тип: {placeholder.name}")
            # print()
            if len(text_predlozenia)<=3 or len(text_forms)==1:
                content_shape = slide.shapes.placeholders[text_forms[0]]
                content_shape.text = ' '.join(text_predlozenia) #slide_content['text']

                for paragraph in content_shape.text_frame.paragraphs:
                    paragraph.font.name = font_name
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.color.rgb = RGBColor(*font_color)

            else:
                count_predl = len(text_predlozenia) // len(text_forms)
                print(count_predl, len(text_predlozenia) ,len(text_forms) )

                if len(text_predlozenia) < len(text_forms):

                    for predloz in range(len(text_predlozenia),):
                        if predloz < len(text_forms)-1:
                            content_shape = slide.shapes.placeholders[text_forms[predloz]]
                            content_shape.text = text_predlozenia[predloz]

                        elif predloz==len(text_forms)-1:
                            content_shape = slide.shapes.placeholders[text_forms[predloz]]
                            content_shape.text = ''.join(text_predlozenia[predloz:])

                        # Настраиваем шрифт основного текста
                        for paragraph in content_shape.text_frame.paragraphs:
                            paragraph.font.name = font_name
                            paragraph.font.size = Pt(font_size)
                            paragraph.font.color.rgb = RGBColor(*font_color)
                            # paragraph.alignment = PP_ALIGN.LEFT
                else:
                    predz = [''.join(text_predlozenia[i:i+count_predl]) for i in range(len(text_predlozenia))]
                    count_form = 0
                    for p in predz:
                        if count_form == len(text_forms)-1:
                            content_shape = slide.shapes.placeholders[text_forms[count_form]]
                            content_shape.text = ''.join(p)
                            for paragraph in content_shape.text_frame.paragraphs:
                                paragraph.font.name = font_name
                                paragraph.font.size = Pt(font_size)
                                paragraph.font.color.rgb = RGBColor(*font_color)
                            break
                        content_shape = slide.shapes.placeholders[text_forms[count_form]]
                        content_shape.text = ''.join(p)

                        for paragraph in content_shape.text_frame.paragraphs:
                          paragraph.font.name = font_name
                          paragraph.font.size = Pt(font_size)
                          paragraph.font.color.rgb = RGBColor(*font_color)

                        count_form += 1




    if design=='1':
        slide_layout = presentation.slide_layouts[5]  # Макет для заголовка и содержимого
        slide = presentation.slides.add_slide(slide_layout)

    slides_to_remove = list(range(0,count_slides_maket))  # Слайды с индексами 0-6 (это шаблон)
    remove_slides(presentation, slides_to_remove)

    presentation.save(output_file)

def generate_presentation_yourself(slides, plan, name, output_file, font_name, font_size, bg_color, font_color="#000000"):
    # Создание новой презентации
    bg_color = hex_to_rgb(bg_color)
    font_color = hex_to_rgb(font_color)

    maket = 'template.pptx'

    presentation = Presentation(maket)

    slides_with_picture = []
    slides_for_text = []

    count_slides_maket = 0
    for slide_number, slide in enumerate(presentation.slides, ):
        count_slides_maket+=1
        flag=True
        for placeholder in slide.placeholders:
            if 'Рисунок' in placeholder.name or 'Picture' in placeholder.name:
                  slides_with_picture.append(slide_number)
                  flag = False
                  break
        if flag:
            slides_for_text.append(slide_number)


    slide_layout = presentation.slide_layouts[0]  # Макет для заголовка и содержимого
    slide = presentation.slides.add_slide(slide_layout)
    slides_with_picture.remove(0)
    # Добавляем заголовок
    title = slide.shapes.title
    title.text = name

    # Настраиваем шрифт заголовка
    title.text_frame.paragraphs[0].font.name = font_name
    title.text_frame.paragraphs[0].font.size = Pt(font_size + 20)
    title.text_frame.paragraphs[0].font.bold = True

    # Устанавливаем цвет фона слайда
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg_color)

    # Проходим по каждому элементу содержания (тематика и текст для слайда)
    for slide_content, title_content in zip(slides, plan):
        
        

        if slide_content['image_path'] != '':
            print(slides_with_picture)
            slide_layout = presentation.slide_layouts[3]  # Макет для заголовка и содержимого
            slide = presentation.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)

            # Добавляем заголовок
            title = slide.shapes.title
            title.text = title_content

            # Настраиваем шрифт заголовка
            title.text_frame.paragraphs[0].font.name = font_name
            title.text_frame.paragraphs[0].font.size = Pt(font_size + 10)
            title.text_frame.paragraphs[0].font.bold = True

            text_forms = []
            img_forms = []

            for placeholder in slide.placeholders:
                if 'Text' in placeholder.name or 'Текст' in placeholder.name:
                    text_forms.append(placeholder.placeholder_format.idx)
                if 'Picture' in placeholder.name or 'Рисунок' in placeholder.name:
                    img_forms.append(placeholder.placeholder_format.idx)

                # print(f"Индекс: {placeholder.placeholder_format.idx}, Тип: {placeholder.name}")

            content_shape = slide.shapes.placeholders[text_forms[0]]
            content_shape.text = slide_content['text']

            content_shape2 = slide.shapes.placeholders[img_forms[0]]
            content_shape2.insert_picture(slide_content['image_path'])

            # Настраиваем шрифт основного текста
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = RGBColor(*font_color)
                # paragraph.alignment = PP_ALIGN.LEFT

        else:

            slide_layout = presentation.slide_layouts[1]  # Макет для заголовка и содержимого
            slide = presentation.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)

            # Добавляем заголовок
            title = slide.shapes.title
            title.text = title_content
                        # Настраиваем шрифт заголовка
            title.text_frame.paragraphs[0].font.name = font_name
            title.text_frame.paragraphs[0].font.size = Pt(font_size + 10)
            title.text_frame.paragraphs[0].font.bold = True

            text_forms = []
            img_forms = []
            
            for placeholder in slide.placeholders:
                if 'Text' in placeholder.name:
                    text_forms.append(placeholder.placeholder_format.idx)
                if 'Picture' in placeholder.name:
                    img_forms.append(placeholder.placeholder_format.idx)

            try:
                text_predlozenia = sent_tokenize(slide_content['text'],)
            except:
                text_predlozenia = list(map(str, slide_content['text'].split('.')))

            if len(text_predlozenia)==1:
                content_shape = slide.shapes.placeholders[text_forms[0]]
                content_shape.text = slide_content['text']

                for paragraph in content_shape.text_frame.paragraphs:
                    paragraph.font.name = font_name
                    paragraph.font.size = Pt(font_size)
                    paragraph.font.color.rgb = RGBColor(*font_color)

            else:
                for predloz in range(len(text_predlozenia)):
                    if predloz < len(text_forms)-1:
                        content_shape = slide.shapes.placeholders[text_forms[predloz]]
                        content_shape.text = text_predlozenia[predloz]

                    elif predloz==len(text_forms)-1:
                        content_shape = slide.shapes.placeholders[text_forms[predloz]]
                        content_shape.text = '. '.join(text_predlozenia[predloz:])
                    # Настраиваем шрифт основного текста
                    for paragraph in content_shape.text_frame.paragraphs:
                        paragraph.font.name = font_name
                        paragraph.font.size = Pt(font_size)
                        paragraph.font.color.rgb = RGBColor(*font_color)
                        # paragraph.alignment = PP_ALIGN.LEFT


    slides_to_remove = list(range(0,count_slides_maket))  # Слайды с индексами 0-6 (это шаблон)
    remove_slides(presentation, slides_to_remove)

    presentation.save(output_file)

# generate_presentation_yourself(slides, plan, name, 'презентация_шаблон_cвой_дизайн.pptx', font_name, font_size, bg_color )